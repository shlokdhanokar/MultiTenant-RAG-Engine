"""
PPTX layout parser.

Slide decks have an inherent hierarchy the pipeline can use directly: each
slide's title placeholder is a Topic and its body text is the content beneath
it, so one slide maps cleanly onto one topic group. Speaker notes are included
because they frequently carry the substance the slide only gestures at.
"""
import logging

logger = logging.getLogger(__name__)


def analyze_pptx_layout(pptx_bytes: bytes):
    """
    Extracts structured blocks and images from a .pptx file.
    page_number is the slide index, which keeps image anchoring per-slide.
    """
    import io
    from pptx import Presentation

    try:
        presentation = Presentation(io.BytesIO(pptx_bytes))
    except Exception as e:
        logger.error(f"failed to load pptx from memory: {e}")
        return None

    semantic_blocks = []
    images = []

    for slide_index, slide in enumerate(presentation.slides):
        title_shape = slide.shapes.title if slide.shapes.title is not None else None
        title_text = (title_shape.text or "").strip() if title_shape is not None else ""
        # python-pptx returns a fresh wrapper object per access, so identity
        # comparison against shapes yielded by iteration never matches — the
        # underlying shape_id is what actually identifies the title placeholder.
        title_shape_id = title_shape.shape_id if title_shape is not None else None

        if title_text:
            semantic_blocks.append({
                "page_number": slide_index,
                "text": title_text,
                "font_size": 0,
                "bbox": (0, 0, 0, 0),  # title sits at the top of its slide
                "type": "Topic",
            })
        else:
            # Untitled slides still need an anchor so their body text isn't
            # absorbed into the previous slide's topic.
            semantic_blocks.append({
                "page_number": slide_index,
                "text": f"Slide {slide_index + 1}",
                "font_size": 0,
                "bbox": (0, 0, 0, 0),
                "type": "Topic",
            })

        for shape in slide.shapes:
            if title_shape_id is not None and shape.shape_id == title_shape_id:
                continue

            if shape.shape_type == 13 or getattr(shape, "image", None) is not None:
                try:
                    image = shape.image
                    images.append({
                        "page_number": slide_index,
                        "image_bytes": image.blob,
                        "image_ext": image.ext,
                        "width": 0,
                        "height": 0,
                        "bbox": (0, float(shape.top or 0), 0, float(shape.top or 0)),
                    })
                except Exception:
                    pass
                continue

            if not shape.has_text_frame:
                continue

            text = (shape.text_frame.text or "").strip()
            if text:
                semantic_blocks.append({
                    "page_number": slide_index,
                    "text": text,
                    "font_size": 0,
                    "bbox": (0, float(shape.top or 1), 0, float(shape.top or 1)),
                    "type": "Paragraph",
                })

        if slide.has_notes_slide:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes:
                semantic_blocks.append({
                    "page_number": slide_index,
                    "text": notes,
                    "font_size": 0,
                    "bbox": (0, 999999, 0, 999999),  # notes sort last within the slide
                    "type": "Paragraph",
                })

    return {
        "total_pages": len(presentation.slides),
        "semantic_blocks": semantic_blocks,
        "images": images,
    }
